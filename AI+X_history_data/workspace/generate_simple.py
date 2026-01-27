from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

# 定义颜色
COLOR_PRIMARY = RGBColor(39, 120, 132)  # #277884
COLOR_ACCENT = RGBColor(254, 68, 71)    # #FE4447
COLOR_SECONDARY = RGBColor(94, 168, 167) # #5EA8A7
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_DARK = RGBColor(51, 51, 51)

# 幻灯片1：标题页
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide1.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = COLOR_PRIMARY

title_box = slide1.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "AI+X 社群聊天记录分析"
title_frame.paragraphs[0].font.size = Pt(48)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.color.rgb = COLOR_WHITE
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(2.8), Inches(8), Inches(0.6))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "社区活动与知识分享回顾"
subtitle_frame.paragraphs[0].font.size = Pt(24)
subtitle_frame.paragraphs[0].font.color.rgb = COLOR_ACCENT
subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

date_box = slide1.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(0.5))
date_frame = date_box.text_frame
date_frame.text = "2025年5月"
date_frame.paragraphs[0].font.size = Pt(18)
date_frame.paragraphs[0].font.color.rgb = COLOR_SECONDARY
date_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# 幻灯片2：社区概况
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
header = slide2.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
header.fill.solid()
header.fill.fore_color.rgb = COLOR_PRIMARY
header.line.color.rgb = COLOR_ACCENT
header.line.width = Pt(5)
header_text = header.text_frame
header_text.text = "社区概况"
header_text.paragraphs[0].font.size = Pt(36)
header_text.paragraphs[0].font.bold = True
header_text.paragraphs[0].font.color.rgb = COLOR_WHITE

content_box = slide2.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(4))
tf = content_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.text = "AI+X 社区核心特色"
p.font.size = Pt(24)
p.font.bold = True
p.font.color.rgb = COLOR_PRIMARY
p.space_after = Pt(15)

features = [
    "非商业化、纯粹的AI知识分享平台",
    "分布式自媒体联动，打造AI知识宇宙",
    "微信知识库为主战场，收录社区成员优质文章",
    "贡献值体系，让每个人都被看到",
    "线下活动组织，华东、华南地区网友见面"
]

for feature in features:
    p = tf.add_paragraph()
    p.text = feature
    p.font.size = Pt(16)
    p.font.color.rgb = COLOR_DARK
    p.level = 0
    p.space_after = Pt(10)

# 幻灯片3：5月活动
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
header = slide3.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
header.fill.solid()
header.fill.fore_color.rgb = COLOR_PRIMARY
header.line.color.rgb = COLOR_ACCENT
header.line.width = Pt(5)
header_text = header.text_frame
header_text.text = "5月重点活动"
header_text.paragraphs[0].font.size = Pt(36)
header_text.paragraphs[0].font.bold = True
header_text.paragraphs[0].font.color.rgb = COLOR_WHITE

# 左侧
left_title = slide3.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4), Inches(0.5))
tf = left_title.text_frame
tf.text = "知识库建设"
tf.paragraphs[0].font.size = Pt(20)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = COLOR_PRIMARY

left_content = slide3.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4), Inches(2))
tf = left_content.text_frame
for item in ["启动微信IMA知识库", "收录社区成员公众号文章", "分门别类整理AI主题", "支持智能问答"]:
    p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
    p.text = item
    p.font.size = Pt(14)
    p.level = 0

highlight1 = slide3.shapes.add_shape(1, Inches(0.5), Inches(3.8), Inches(4), Inches(0.8))
highlight1.fill.solid()
highlight1.fill.fore_color.rgb = COLOR_SECONDARY
tf = highlight1.text_frame
tf.text = "163篇文章已收录"
tf.paragraphs[0].font.size = Pt(18)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = COLOR_WHITE
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# 右侧
right_title = slide3.shapes.add_textbox(Inches(5.5), Inches(1.2), Inches(4), Inches(0.5))
tf = right_title.text_frame
tf.text = "社区贡献激励"
tf.paragraphs[0].font.size = Pt(20)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = COLOR_PRIMARY

right_content = slide3.shapes.add_textbox(Inches(5.5), Inches(1.8), Inches(4), Inches(2))
tf = right_content.text_frame
for item in ["4月份贡献度统计", "最佳问题奖：佑禾", "最佳创意奖：御风现梦", "最佳工具奖：张三"]:
    p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
    p.text = item
    p.font.size = Pt(14)
    p.level = 0

highlight2 = slide3.shapes.add_shape(1, Inches(5.5), Inches(3.8), Inches(4), Inches(0.8))
highlight2.fill.solid()
highlight2.fill.fore_color.rgb = COLOR_SECONDARY
tf = highlight2.text_frame
tf.text = "每月评选\n让贡献者被看到"
tf.paragraphs[0].font.size = Pt(16)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = COLOR_WHITE
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# 幻灯片4：工具分享
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
header = slide4.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
header.fill.solid()
header.fill.fore_color.rgb = COLOR_PRIMARY
header.line.color.rgb = COLOR_ACCENT
header.line.width = Pt(5)
header_text = header.text_frame
header_text.text = "技术工具分享"
header_text.paragraphs[0].font.size = Pt(36)
header_text.paragraphs[0].font.bold = True
header_text.paragraphs[0].font.color.rgb = COLOR_WHITE

title_box = slide4.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.5))
tf = title_box.text_frame
tf.text = "社区成员贡献的实用工具"
tf.paragraphs[0].font.size = Pt(20)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = COLOR_PRIMARY

tools = [
    ("公众号导出", "张三分享的工具，可一键导出公众号文章"),
    ("AI视频生成", "御风现梦研究的多种风格视频生成技术"),
    ("表情包制作", "探索AI生成GIF表情包的流程化方案")
]

for i, (title, desc) in enumerate(tools):
    x = 0.5 + i * 3.2
    card = slide4.shapes.add_shape(1, Inches(x), Inches(2), Inches(3), Inches(1.5))
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(245, 245, 245)
    card.line.color.rgb = COLOR_SECONDARY
    card.line.width = Pt(3)

    tf = card.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(12)
    p2.font.color.rgb = COLOR_DARK

quote = slide4.shapes.add_shape(1, Inches(1), Inches(4), Inches(8), Inches(0.8))
quote.fill.solid()
quote.fill.fore_color.rgb = COLOR_ACCENT
tf = quote.text_frame
tf.text = '"开始用数据资产的视角去写和保存公众号"'
tf.paragraphs[0].font.size = Pt(16)
tf.paragraphs[0].font.italic = True
tf.paragraphs[0].font.color.rgb = COLOR_WHITE
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# 幻灯片5：未来规划
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
header = slide5.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
header.fill.solid()
header.fill.fore_color.rgb = COLOR_PRIMARY
header.line.color.rgb = COLOR_ACCENT
header.line.width = Pt(5)
header_text = header.text_frame
header_text.text = "未来规划"
header_text.paragraphs[0].font.size = Pt(36)
header_text.paragraphs[0].font.bold = True
header_text.paragraphs[0].font.color.rgb = COLOR_WHITE

title_box = slide5.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.5))
tf = title_box.text_frame
tf.text = "6月活动预告"
tf.paragraphs[0].font.size = Pt(20)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = COLOR_PRIMARY

events = [
    ("6月", "组织华东地区小伙伴前往华南进行线下交流"),
    ("持续", "继续完善知识库内容，收录更多优质AI文章"),
    ("探索", "体验AI+X小程序，探索更多AI应用场景"),
    ("合作", "带着社区知识库寻求与腾讯等平台的合作")
]

y = 2
for badge, desc in events:
    badge_shape = slide5.shapes.add_shape(1, Inches(0.5), Inches(y), Inches(1), Inches(0.4))
    badge_shape.fill.solid()
    badge_shape.fill.fore_color.rgb = COLOR_ACCENT
    tf = badge_shape.text_frame
    tf.text = badge
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = COLOR_WHITE
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    content_shape = slide5.shapes.add_shape(1, Inches(1.8), Inches(y), Inches(7.7), Inches(0.4))
    content_shape.fill.solid()
    content_shape.fill.fore_color.rgb = RGBColor(245, 245, 245)
    content_shape.line.color.rgb = COLOR_SECONDARY
    content_shape.line.width = Pt(2)
    tf = content_shape.text_frame
    tf.text = desc
    tf.paragraphs[0].font.size = Pt(13)
    tf.paragraphs[0].font.color.rgb = COLOR_DARK

    y += 0.7

# 幻灯片6：愿景
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
background = slide6.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = COLOR_PRIMARY

quote_box = slide6.shapes.add_shape(1, Inches(1.5), Inches(1.5), Inches(7), Inches(2))
quote_box.fill.solid()
quote_box.fill.fore_color.rgb = RGBColor(255, 255, 255)
quote_box.fill.transparency = 0.9
quote_box.line.color.rgb = COLOR_ACCENT
quote_box.line.width = Pt(3)

tf = quote_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '"提供真正的价值，穿越时间周期波动，做最共享普惠的 AI + X 分享"'
p.font.size = Pt(24)
p.font.color.rgb = COLOR_WHITE
p.alignment = PP_ALIGN.CENTER
p.space_after = Pt(15)

p2 = tf.add_paragraph()
p2.text = "—— AI+X 社区愿景"
p2.font.size = Pt(16)
p2.font.color.rgb = COLOR_SECONDARY
p2.alignment = PP_ALIGN.CENTER

vision_box = slide6.shapes.add_textbox(Inches(2), Inches(4), Inches(6), Inches(0.6))
tf = vision_box.text_frame
tf.text = "给所有纯粹的AI文章们一个家"
tf.paragraphs[0].font.size = Pt(20)
tf.paragraphs[0].font.color.rgb = COLOR_ACCENT
tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# 保存
prs.save('D:/AI+X web/AI+X_history_data/AI+X_社群分析_2025年5月.pptx')
print("PPT生成成功！")
